"""
Build "GenbAi" executive pitch deck as a PowerPoint file.
GenbAi = Genba (Toyota: go to the real place) + AI (silent i).
Framing: "Art of the Possible" — built in a single day with Claude Code.
Target audience: Senior executives, not investors.
Palette: navy / blue / steel — matched to GenbAi logo.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

OUT_PATH  = "/Users/pwrhome/Claude Code Working/the-cascade/GenbAi_ExecutiveDeck.pptx"
LOGO_PATH = "/Users/pwrhome/Claude Code Working/the-cascade/frontend/public/genbai-logo.png"

# ── Brand palette (GenbAi logo-matched) ─────────────────────────────────────
NAVY        = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy  — dark slide backgrounds
BLACK       = RGBColor(0x11, 0x11, 0x11)   # near-black — body text on light slides
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLUE        = RGBColor(0x1E, 0x7F, 0xD0)   # logo blue  — primary brand accent
STEEL       = RGBColor(0x94, 0xA3, 0xB8)   # silver-steel — secondary accent
RED         = RGBColor(0xCC, 0x22, 0x22)   # pulse red  — alerts / problem stats
LIGHT_GRAY  = RGBColor(0xF4, 0xF4, 0xF5)   # near-white — light-bg slides
MID_GRAY    = RGBColor(0x88, 0x8C, 0x99)   # cool mid-gray — supporting text
DARK_SLATE  = RGBColor(0x1E, 0x2A, 0x38)   # dark slate — panel fill on dark slides
DARK_GRAY   = RGBColor(0x2C, 0x35, 0x45)   # mid-dark   — alternate row / panel fill

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank


# ── helpers ──────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(blank_layout)

def rect(slide, x, y, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(
        1,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def textbox(slide, x, y, w, h, text, size=18, bold=False, color=BLACK,
            align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def multiline(slide, x, y, w, h, lines, size=14, color=BLACK, bold_first=False, spacing=None):
    """lines = list of strings. First line optionally bold."""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if spacing:
            p.space_before = Pt(spacing)
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = (bold_first and i == 0)
        run.font.color.rgb = color
    return txb

def accent_bar(slide, color=BLUE):
    rect(slide, 0, 0, 13.33, 0.08, fill=color)

def add_logo(slide, x, y, w):
    """Embed GenbAi logo if the file exists at LOGO_PATH."""
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(x), Inches(y), width=Inches(w))
        return True
    return False


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=NAVY)
rect(s, 0, 0, 13.33, 0.12, fill=BLUE)
rect(s, 0, 7.38, 13.33, 0.12, fill=BLUE)

logo_shown = add_logo(s, 4.67, 0.55, 4.0)   # centered logo ~4" wide

if logo_shown:
    # Logo is present — it includes the wordmark; place tagline below
    textbox(s, 1, 3.35, 11, 0.6,
            "Go to the genba. Every time.",
            size=26, color=BLUE, align=PP_ALIGN.CENTER, italic=True)
    textbox(s, 1, 4.1, 11, 0.5,
            "AI-Powered Crisis Simulation  ·  Built in One Day  ·  Digital Innovations",
            size=15, color=MID_GRAY, align=PP_ALIGN.CENTER)
    textbox(s, 1, 4.8, 11, 0.4,
            "Scenario: The Cascade",
            size=13, color=STEEL, align=PP_ALIGN.CENTER)
else:
    # Text-only fallback
    textbox(s, 1, 1.6, 11, 1.4,
            "GenbAi",
            size=88, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(s, 1, 3.1, 11, 0.55,
            "Go to the genba. Every time.",
            size=26, color=BLUE, align=PP_ALIGN.CENTER, italic=True)
    textbox(s, 1, 3.8, 11, 0.5,
            "AI-Powered Crisis Simulation  ·  Built in One Day  ·  Digital Innovations",
            size=15, color=MID_GRAY, align=PP_ALIGN.CENTER)
    textbox(s, 1, 4.55, 11, 0.4,
            "Scenario: The Cascade",
            size=13, color=STEEL, align=PP_ALIGN.CENTER)

textbox(s, 1, 6.5, 11, 0.5,
        "Confidential  ·  Executive Preview  ·  2026",
        size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE NAME SAYS EVERYTHING
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=NAVY)
accent_bar(s)

textbox(s, 0.5, 0.3, 12, 0.6,
        "The Name Says Everything", size=28, bold=True, color=WHITE)

# Left — Genba (Toyota principle)
rect(s, 0.4, 1.1, 5.9, 5.5, fill=DARK_SLATE)
rect(s, 0.4, 1.1, 5.9, 0.1, fill=BLUE)
textbox(s, 0.55, 1.3, 5.6, 0.55, "現場  ·  Genba", size=22, bold=True, color=BLUE)
textbox(s, 0.55, 2.0, 5.6, 0.4, "The Toyota Principle", size=13, color=MID_GRAY, italic=True)
multiline(s, 0.55, 2.55, 5.6, 3.8, [
    "In the Toyota Production System, leaders",
    "don't manage from a boardroom.",
    "",
    "They go to the genba —",
    "the real place, where the work happens —",
    "to see, understand, and improve.",
    "",
    "You cannot solve a problem",
    "you haven't stood inside.",
], size=13, color=WHITE)

# Right — GenbAi (Digital Innovations platform)
rect(s, 6.9, 1.1, 6.0, 5.5, fill=DARK_SLATE)
rect(s, 6.9, 1.1, 6.0, 0.1, fill=STEEL)
textbox(s, 7.05, 1.3, 5.7, 0.55, "GenbAi", size=32, bold=True, color=WHITE)
textbox(s, 7.05, 2.0, 5.7, 0.4, "The Digital Innovations Platform", size=13, color=STEEL, italic=True)
multiline(s, 7.05, 2.55, 5.7, 3.8, [
    "GenbAi brings the genba principle",
    "into the age of artificial intelligence.",
    "",
    "Players don't sit through a slide deck.",
    "They step inside a live crisis —",
    "make real decisions, under real pressure,",
    "with AI as their instrument.",
    "",
    "The 'i' is silent. The AI is not.",
], size=13, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE PROBLEM: AI Upskilling at Scale
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)
accent_bar(s)

textbox(s, 0.5, 0.3, 10, 0.6,
        "The Problem: AI Upskilling at Scale",
        size=28, bold=True, color=BLACK)

# Three stat boxes
for i, (stat, label, col) in enumerate([
    ("73%",   "of employees feel underprepared\nfor AI-augmented roles",          RED),
    ("< 5%",  "of AI training programmes\nshow sustained behaviour change",       STEEL),
    ("18 mo", "average time to deploy\nenterprise-wide AI curriculum",            DARK_GRAY),
]):
    bx = 0.4 + i * 4.3
    rect(s, bx, 1.2, 3.9, 2.8, fill=WHITE, line=col)
    rect(s, bx, 1.2, 3.9, 0.08, fill=col)
    textbox(s, bx + 0.15, 1.4, 3.6, 1.0, stat, size=52, bold=True, color=col)
    multiline(s, bx + 0.15, 2.5, 3.6, 1.2, label.split('\n'), size=13, color=DARK_GRAY)

textbox(s, 0.5, 4.3, 12, 0.5,
        "The traditional answer — classroom training — is slow, expensive, and doesn't scale.",
        size=15, color=DARK_GRAY, italic=True)

multiline(s, 0.5, 4.9, 12, 2.2, [
    "•  Cohort-based workshops serve hundreds; organisations need to reach thousands.",
    "•  Engagement drops sharply after the first session — retention falls to < 20% within 30 days.",
    "•  Generic scenarios don't connect to employees' actual work, reducing transfer of learning.",
    "•  Facilitator costs prevent continuous, repeatable deployment across geographies.",
], size=13, color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — TRADITIONAL TRAINING IS FAILING
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=WHITE)
accent_bar(s, RED)

textbox(s, 0.5, 0.3, 12, 0.6,
        "Why Traditional AI Training Falls Short", size=28, bold=True, color=BLACK)

rows = [
    ("Dimension",        "Classroom / eLearning",         "What's Needed",            True),
    ("Scale",            "Hundreds per quarter",          "Thousands simultaneously", False),
    ("Engagement",       "Passive, slide-driven",         "Active, high-stakes",      False),
    ("Relevance",        "Generic case studies",          "Real business scenarios",  False),
    ("Retention",        "< 20% after 30 days",          "> 60% applied behaviour",  False),
    ("Cost per learner", "High (facilitator + venue)",    "Marginal cost near zero",  False),
    ("Feedback loop",    "Quarterly survey",              "Real-time analytics",      False),
]

col_x = [0.4, 3.5, 8.1]
col_w = [2.9, 4.4, 4.8]
row_h = 0.6
start_y = 1.1

for r_i, row in enumerate(rows):
    bg = LIGHT_GRAY if r_i % 2 == 0 else WHITE
    if r_i == 0:
        bg = NAVY
    rect(s, 0.4, start_y + r_i * row_h, 12.5, row_h, fill=bg)
    for c_i, cell in enumerate(row[:3]):
        col = WHITE if r_i == 0 else (RED if (c_i == 1 and r_i > 0) else (BLUE if c_i == 2 and r_i > 0 else BLACK))
        b = (r_i == 0)
        textbox(s, col_x[c_i] + 0.1, start_y + r_i * row_h + 0.1,
                col_w[c_i] - 0.1, row_h - 0.1,
                cell, size=12, bold=b, color=col)

textbox(s, 0.5, 6.6, 12, 0.5,
        "The gap between what training delivers and what transformation requires is widening — not closing.",
        size=12, italic=True, color=MID_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — INTRODUCING GenbAi
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=NAVY)
accent_bar(s)

textbox(s, 0.5, 0.3, 12, 0.6,
        "Introducing: GenbAi", size=32, bold=True, color=WHITE)
textbox(s, 0.5, 1.0, 12, 0.4,
        "An immersive, AI-driven crisis simulation that upskills and engages simultaneously.",
        size=16, color=BLUE, italic=True)

for i, (icon, title, body) in enumerate([
    ("⚡", "Real Stakes",       "48-hour countdown to a production line stoppage. Players feel the pressure — no passive observation."),
    ("🤝", "Dual Role Play",    "Two personas. Two data sets. Neither has the full picture. Cross-functional collaboration is mandatory."),
    ("🧠", "AI as a Tool",      "Players use AI to interrogate documents, surface anomalies, and synthesise insights under time pressure."),
    ("📊", "Real-World Scenario","Modelled on actual supply chain risks: supplier insolvency, quality holds, logistics failures."),
]):
    bx = 0.5 + i * 3.1
    rect(s, bx, 1.8, 2.8, 3.2, fill=DARK_SLATE)
    textbox(s, bx + 0.1, 1.9, 2.6, 0.5, icon, size=28, color=WHITE)
    textbox(s, bx + 0.1, 2.5, 2.6, 0.4, title, size=14, bold=True, color=BLUE)
    multiline(s, bx + 0.1, 3.0, 2.6, 1.8, body.split('. '), size=11, color=WHITE)

textbox(s, 0.5, 5.3, 12, 0.5,
        "Not a training module. Not a workshop. A live simulation your teams will actually want to complete.",
        size=14, color=STEEL, italic=True, align=PP_ALIGN.CENTER)

rect(s, 0.5, 5.9, 12, 1.2, fill=DARK_GRAY)
textbox(s, 0.5, 6.0, 12, 0.5,
        "Powered by Claude (Anthropic)  ·  Real-time AI coaching  ·  Adaptive scenario scoring  ·  WebSocket multiplayer",
        size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — DUAL LEARNING OBJECTIVES
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=WHITE)
accent_bar(s)

textbox(s, 0.5, 0.3, 12, 0.6,
        "One Simulation. Two Learning Objectives.", size=28, bold=True, color=BLACK)

# Left — AI / Digital Fluency
rect(s, 0.4, 1.1, 5.9, 5.8, fill=RGBColor(0xE6, 0xF2, 0xFB))   # light blue
rect(s, 0.4, 1.1, 5.9, 0.1, fill=BLUE)
textbox(s, 0.5, 1.25, 5.7, 0.5, "① AI / Digital Fluency", size=18, bold=True, color=BLUE)
multiline(s, 0.5, 1.85, 5.7, 4.8, [
    "Prompt engineering for document analysis",
    "AI-assisted hypothesis formation",
    "Cross-referencing AI outputs vs raw data",
    "Recognising AI limitations and hallucination risk",
    "Iterative querying to refine insights",
    "Using AI to synthesise multi-source intelligence",
    "Speed vs accuracy trade-offs under pressure",
    "Trust calibration: when to override AI suggestions",
], size=13, color=BLACK)

# Right — Business / Domain Mastery
rect(s, 6.8, 1.1, 6.1, 5.8, fill=RGBColor(0xF0, 0xF4, 0xF8))   # light steel
rect(s, 6.8, 1.1, 6.1, 0.1, fill=STEEL)
textbox(s, 6.9, 1.25, 5.9, 0.5, "② Business / Domain Mastery", size=18, bold=True, color=DARK_GRAY)
multiline(s, 6.9, 1.85, 5.9, 4.8, [
    "Supply chain risk identification and triage",
    "Vendor financial health analysis",
    "Cross-functional crisis communication",
    "Root cause analysis under incomplete data",
    "Escalation judgment and decision authority",
    "Trade-off analysis: cost vs continuity",
    "Stakeholder alignment in time-pressured decisions",
    "Structured problem-solving frameworks",
], size=13, color=BLACK)

textbox(s, 0.5, 7.1, 12, 0.3,
        "Most training programmes pick one. GenbAi delivers both — simultaneously, in context.",
        size=11, italic=True, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — HOW IT WORKS (GAMEPLAY)
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)
accent_bar(s)

textbox(s, 0.5, 0.3, 12, 0.6,
        "How It Works: The Player Experience", size=28, bold=True, color=BLACK)

steps = [
    ("1", "JOIN",        "Each player enters a unique session link and selects their role — Supply Chain Manager or Finance Analyst.", BLUE),
    ("2", "INVESTIGATE", "Players explore a dossier of 8 role-specific documents: emails, spreadsheets, alerts, financial reports.",  STEEL),
    ("3", "COLLABORATE", "Real-time War Room chat. Neither player can solve the crisis alone — coordination is essential.",            RED),
    ("4", "DECIDE",      "Teams formulate a root cause and recommend a corrective action under the countdown.",                      DARK_GRAY),
    ("5", "DEBRIEF",     "AI evaluates submissions, scores the response, and delivers a personalised debrief with improvement coaching.", BLUE),
]

for i, (num, title, body, col) in enumerate(steps):
    y = 1.1 + i * 1.2
    rect(s, 0.4, y, 0.5, 1.0, fill=col)
    textbox(s, 0.4, y + 0.15, 0.5, 0.7, num, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(s, 1.1, y + 0.05, 1.8, 0.4, title, size=14, bold=True, color=col)
    textbox(s, 1.1, y + 0.45, 11.5, 0.5, body, size=12, color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — TECHNOLOGY STACK
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=NAVY)
accent_bar(s)

textbox(s, 0.5, 0.3, 12, 0.6,
        "Technology: What's Under the Hood", size=28, bold=True, color=WHITE)

tech = [
    ("Frontend",         "Next.js 14, TypeScript, Tailwind CSS, React Markdown",       "Real-time reactive UI"),
    ("Backend",          "Python FastAPI, WebSocket, Redis",                            "Persistent sessions, sub-100ms latency"),
    ("AI — Coaching",   "Claude Haiku  (Anthropic)",                                   "Adaptive hints during play"),
    ("AI — Evaluation", "Claude Opus  (Anthropic)",                                    "Intelligent scoring + personalised debrief"),
    ("AI — Scenario",   "Claude Sonnet  (Anthropic)",                                  "Scenario generation and content creation"),
    ("Infra",            "Stateless API, Redis TTL sessions, environment-configurable", "Cloud-ready, multi-tenant capable"),
]

for i, (layer, tech_str, note) in enumerate(tech):
    y = 1.2 + i * 0.95
    bg = DARK_SLATE if i % 2 == 0 else DARK_GRAY
    rect(s, 0.4, y, 12.5, 0.85, fill=bg)
    textbox(s, 0.55, y + 0.1, 2.2, 0.65, layer, size=12, bold=True, color=BLUE)
    textbox(s, 2.9,  y + 0.1, 5.5, 0.65, tech_str, size=12, color=WHITE)
    textbox(s, 8.6,  y + 0.1, 4.2, 0.65, note, size=11, color=MID_GRAY, italic=True)

textbox(s, 0.5, 7.0, 12, 0.35,
        "Fully browser-based. No app install. No VPN. Runs anywhere.",
        size=12, color=STEEL, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — BUILT IN ONE DAY
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=WHITE)
accent_bar(s, STEEL)

textbox(s, 0.5, 0.3, 12, 0.6,
        "Built in One Day — With Claude Code", size=28, bold=True, color=BLACK)
textbox(s, 0.5, 1.0, 12, 0.4,
        "From blank repository to fully playable, AI-coached multiplayer simulation in a single working session.",
        size=14, color=STEEL, italic=True)

timeline = [
    ("Morning",   "Concept & Architecture", "Scenario design, persona definition, document dossier creation, API contracts sketched"),
    ("Mid-Day",   "Core Build",             "FastAPI backend, WebSocket real-time layer, Redis sessions, game state machine"),
    ("Afternoon", "AI Integration",         "Claude Haiku coaching agent, Opus evaluator, adaptive hint triggers, scoring rubric"),
    ("Late Day",  "Frontend",               "Next.js UI, document viewer, War Room chat, countdown timer, markdown rendering"),
    ("Evening",   "Polish & Refinement",    "Brand reskin, scenario content upgrade, gate questions, escape submissions, debrief flow"),
]

for i, (time, phase, detail) in enumerate(timeline):
    y = 1.6 + i * 1.0
    rect(s, 0.4, y, 1.5, 0.8, fill=BLUE)
    textbox(s, 0.4, y + 0.1, 1.5, 0.65, time, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, 0.4, y + 0.38, 12.5, 0.02, fill=RGBColor(0xDD, 0xDD, 0xDD))
    textbox(s, 2.1, y + 0.02, 3.0, 0.4, phase, size=13, bold=True, color=BLACK)
    textbox(s, 2.1, y + 0.42, 10.8, 0.4, detail, size=11, color=MID_GRAY)

textbox(s, 0.5, 6.8, 12, 0.4,
        "Every architectural decision, line of code, UI component, and AI prompt crafted through conversational AI collaboration.",
        size=11, italic=True, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — POC vs FULL BUILD
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)
accent_bar(s)

textbox(s, 0.5, 0.3, 12, 0.6,
        "POC Today → Full Platform Tomorrow", size=28, bold=True, color=BLACK)

# POC column
rect(s, 0.4, 1.1, 5.9, 6.0, fill=WHITE, line=STEEL)
rect(s, 0.4, 1.1, 5.9, 0.55, fill=STEEL)
textbox(s, 0.5, 1.15, 5.7, 0.4, "✓  Proof of Concept — What Exists Today", size=14, bold=True, color=WHITE)
poc_items = [
    "Single scenario (supply chain crisis)",
    "2-player collaborative session",
    "8 documents per persona",
    "AI coaching hints on request",
    "AI-powered evaluation + debrief",
    "Real-time War Room chat",
    "Countdown timer + gate questions",
    "Session persistence via Redis",
    "Fully browser-based, no install",
    "Hardcoded session / no auth layer",
]
multiline(s, 0.5, 1.8, 5.7, 5.0, ["• " + it for it in poc_items], size=12, color=DARK_GRAY)

# Full build column
rect(s, 6.9, 1.1, 6.0, 6.0, fill=WHITE, line=BLUE)
rect(s, 6.9, 1.1, 6.0, 0.55, fill=BLUE)
textbox(s, 7.0, 1.15, 5.8, 0.4, "★  Full Platform — What a Build Unlocks", size=14, bold=True, color=WHITE)
full_items = [
    "Scenario library (Finance, HR, Legal, Ops…)",
    "N-player sessions (teams of 4–10)",
    "SSO / corporate identity integration",
    "Facilitator dashboard & live monitoring",
    "Analytics: engagement, performance, gaps",
    "Scenario authoring tool for L&D teams",
    "Adaptive difficulty based on player behaviour",
    "Multi-language / multi-market support",
    "LMS integration (Workday, Cornerstone)",
    "Compliance and certification pathways",
]
multiline(s, 7.0, 1.8, 5.8, 5.0, ["• " + it for it in full_items], size=12, color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — AI SKILLS THE SIMULATION TEACHES
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=NAVY)
accent_bar(s)

textbox(s, 0.5, 0.3, 12, 0.6,
        "AI Skills Participants Build", size=28, bold=True, color=WHITE)

skills = [
    ("Prompt Engineering",         "Craft precise questions to extract specific insights from AI across documents, alerts, and reports."),
    ("Hypothesis Formation",       "Use AI to generate and test multiple root-cause theories quickly before committing to a direction."),
    ("Multi-Source Synthesis",     "Direct AI to connect signals across financial data, logistics alerts, supplier emails and quality holds."),
    ("Trust Calibration",          "Evaluate AI outputs critically — knowing when to push back, verify, or override a suggestion."),
    ("Iterative Refinement",       "Progressively sharpen queries based on what AI returns; practise the agentic loop under real time pressure."),
    ("Output Structuring",         "Use AI to format and prioritise findings for a structured decision submission — clarity under pressure."),
    ("Risk & Hallucination Aware.", "Identify when AI response confidence exceeds actual data support; apply healthy scepticism."),
    ("Speed vs Accuracy Balance",  "Develop judgment about when AI-assisted analysis is 'good enough' vs when more rigour is needed."),
]

for i, (skill, desc) in enumerate(skills):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 6.5
    y = 1.1 + row * 1.45
    rect(s, x, y, 6.1, 1.3, fill=DARK_SLATE)
    rect(s, x, y, 6.1, 0.08, fill=BLUE)
    textbox(s, x + 0.15, y + 0.12, 5.8, 0.4, skill, size=13, bold=True, color=BLUE)
    textbox(s, x + 0.15, y + 0.55, 5.8, 0.65, desc, size=11, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — BUSINESS SKILLS THE SCENARIO EXERCISES
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=WHITE)
accent_bar(s, STEEL)

textbox(s, 0.5, 0.3, 12, 0.6,
        "Business Skills the Crisis Demands", size=28, bold=True, color=BLACK)

biz_skills = [
    ("Supply Chain Risk Assessment",    "Identify cascading failure points across supplier, logistics, and quality dimensions simultaneously."),
    ("Financial Health Analysis",       "Interpret credit alerts, overdue receivables, and forex exposure to assess vendor viability."),
    ("Cross-Functional Collaboration",  "Align two roles — operations and finance — with different data sets toward a shared decision."),
    ("Root Cause Analysis",             "Move beyond symptoms to systemic causes using structured deductive reasoning under time pressure."),
    ("Escalation Judgment",             "Determine when a risk warrants senior involvement vs can be contained at team level."),
    ("Trade-off Decision Making",       "Balance cost, continuity, and relationship risk when selecting alternate suppliers or actions."),
    ("Communication Under Pressure",   "Share findings concisely in a live War Room — prioritising signal over noise."),
    ("Structured Problem Solving",      "Apply frameworks (5-Why, Ishikawa) within an AI-assisted context to reach a defensible conclusion."),
]

for i, (skill, desc) in enumerate(biz_skills):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 6.5
    y = 1.1 + row * 1.45
    rect(s, x, y, 6.1, 1.3, fill=LIGHT_GRAY)
    rect(s, x, y, 6.1, 0.08, fill=STEEL)
    textbox(s, x + 0.15, y + 0.12, 5.8, 0.4, skill, size=13, bold=True, color=DARK_GRAY)
    textbox(s, x + 0.15, y + 0.55, 5.8, 0.65, desc, size=11, color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — THE AI COACHING ENGINE
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=NAVY)
accent_bar(s)

textbox(s, 0.5, 0.3, 12, 0.6,
        "Built-In AI Coaching — Not an Add-On", size=28, bold=True, color=WHITE)
textbox(s, 0.5, 1.1, 12, 0.45,
        "Three distinct AI models serve different roles in the experience:",
        size=14, color=MID_GRAY)

coaching_points = [
    (BLUE,     "HAIKU — Real-Time Coach",
               "Monitors player behaviour: files opened, clues found, War Room activity. "
               "Fires contextual hints when teams stall — guiding without giving away the answer. "
               "Responds within seconds, maintaining simulation immersion."),
    (STEEL,    "OPUS — Evaluation & Debrief",
               "Receives both players' final submissions and evaluates root cause accuracy, "
               "action quality, and collaboration evidence. Produces a scored, narrative debrief "
               "that is personal, specific, and coaching-forward."),
    (MID_GRAY, "SONNET — Scenario Intelligence",
               "Used during content creation to generate realistic documents, red herrings, "
               "financial data anomalies, and structured gate questions that test comprehension "
               "at each clue milestone."),
]

for i, (col, title, body) in enumerate(coaching_points):
    y = 1.75 + i * 1.75
    rect(s, 0.4, y, 12.5, 1.6, fill=DARK_SLATE)
    rect(s, 0.4, y, 0.12, 1.6, fill=col)
    textbox(s, 0.7, y + 0.1, 3.5, 0.5, title, size=14, bold=True, color=col)
    textbox(s, 0.7, y + 0.65, 11.8, 0.85, body, size=12, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — WHAT MAKES THIS DIFFERENT
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=WHITE)
accent_bar(s)

textbox(s, 0.5, 0.3, 12, 0.6,
        "Why This Is Different", size=28, bold=True, color=BLACK)

differentiators = [
    ("Context-Native Learning",
     "The scenario mirrors real business complexity — messy data, incomplete information, time pressure. "
     "Learning happens inside the actual context it needs to be applied, not in abstraction."),
    ("AI Is the Tool, Not the Teacher",
     "Players use AI as a working instrument — not as a subject matter expert lecturing at them. "
     "This builds genuine, transferable AI fluency that survives beyond the simulation."),
    ("Every Run Is Unique",
     "Dynamic AI coaching adapts to each team's behaviour. Two groups running the same scenario "
     "will experience meaningfully different hints, pacing, and debrief outcomes."),
    ("Scales Instantly",
     "Once built, the simulation can be deployed to thousands of concurrent sessions with zero "
     "marginal facilitator cost. Global rollout becomes a configuration change, not a project."),
    ("Measurable Output",
     "Evaluation scores, collaboration data, time-to-discovery metrics, and AI usage patterns "
     "create a rich analytics layer for L&D teams that no classroom can provide."),
]

for i, (title, body) in enumerate(differentiators):
    y = 1.1 + i * 1.2
    rect(s, 0.4, y, 0.08, 1.0, fill=BLUE)
    textbox(s, 0.7, y + 0.05, 3.0, 0.45, title, size=13, bold=True, color=BLUE)
    textbox(s, 0.7, y + 0.52, 12.0, 0.5, body, size=12, color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — WHAT NEXT
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)
accent_bar(s)

textbox(s, 0.5, 0.3, 12, 0.6,
        "What Happens Next", size=28, bold=True, color=BLACK)

next_steps = [
    ("Play It",          "Run the POC live with a small cohort — two volunteers, 45 minutes, no setup required."),
    ("Evaluate",         "Collect participant reactions, AI debrief scores, and observed behaviour change against baseline."),
    ("Define Scenarios", "Select 2–3 business-critical domains where AI upskilling + domain practice is highest value."),
    ("Commission Build", "Scope a full GenbAi platform build: scenario authoring, analytics dashboard, SSO, LMS integration."),
    ("Scale",            "Roll out to full target population — same simulation engine, unlimited concurrent sessions."),
]

for i, (step, body) in enumerate(next_steps):
    y = 1.2 + i * 1.1
    rect(s, 0.4, y, 12.5, 1.0, fill=WHITE)
    rect(s, 0.4, y, 0.7,  1.0, fill=BLUE)
    textbox(s, 0.4, y + 0.2, 0.7, 0.6,
            str(i + 1), size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(s, 1.3, y + 0.08, 2.4, 0.45, step, size=14, bold=True, color=BLUE)
    textbox(s, 1.3, y + 0.53, 11.3, 0.4, body, size=12, color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — CLOSING
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=NAVY)
rect(s, 0, 0, 13.33, 0.12, fill=BLUE)
rect(s, 0, 7.38, 13.33, 0.12, fill=BLUE)

textbox(s, 1, 1.4, 11, 1.0,
        "If we built this in one day —",
        size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

textbox(s, 1, 2.6, 11, 0.8,
        "imagine what a dedicated team builds in a quarter.",
        size=32, color=BLUE, align=PP_ALIGN.CENTER, italic=True)

textbox(s, 1, 3.9, 11, 0.6,
        "The only constraint on AI-powered learning at scale is imagination.",
        size=18, color=MID_GRAY, align=PP_ALIGN.CENTER)

textbox(s, 1, 5.1, 11, 0.5,
        "Let's play. Then let's build.",
        size=22, bold=True, color=STEEL, align=PP_ALIGN.CENTER)

# Small logo in bottom-right
add_logo(s, 10.5, 5.9, 2.3)

textbox(s, 1, 6.5, 11, 0.5,
        "GenbAi  ·  Confidential  ·  2026",
        size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(OUT_PATH)
print(f"✅ Saved: {OUT_PATH}")
print(f"   Slides: {len(prs.slides)}")
if not os.path.exists(LOGO_PATH):
    print(f"\n⚠️  Logo not found at: {LOGO_PATH}")
    print("   Save genbai-logo.png there and re-run to embed it on Cover + Closing slides.")
